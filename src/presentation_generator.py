import os
import tkinter as tk
from tkinter import messagebox
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Pt
from datetime import date

class PresentationGenerator:
    def __init__(self, data, file_name, background_image_path):
        self.data = data
        self.file_name = file_name
        self.background_image_path = background_image_path

        # Obtém o caminho de destino das apresentações
        self.dest_path = self.get_dest_path()

    def get_dest_path(self):
        # Tenta ler o caminho de destino do arquivo path.txt
        try:
            with open("path.txt", "r") as file:
                dest_path = file.read().strip()
                if os.path.isdir(dest_path):
                    return dest_path
                else:
                    raise FileNotFoundError
        except FileNotFoundError:
            # Se o arquivo não existir ou o caminho não for válido, pede ao usuário que escolha um novo destino
            dest_path = filedialog.askdirectory(title="Escolha a pasta de destino")
            if dest_path:
                # Salva o novo caminho de destino no arquivo path.txt
                with open("path.txt", "w") as file:
                    file.write(dest_path)
                return dest_path
            else:
                # Caso o usuário cancele a seleção, retorna None
                return None

    def generate_presentation(self):
        # Verifica se o caminho de destino está disponível
        if not self.dest_path:
            messagebox.showerror("Erro", "A pasta de destino não foi selecionada. Por favor, escolha um destino válido.")
            return

        prs = Presentation()

        # Define fonte e tamanho padrão para o texto
        def set_font(text, bold=False, size=25):
            run = text.text_frame.add_paragraph().runs[0]
            font = run.font
            font.bold = bold
            font.name = "Calibri"
            font.size = Pt(size)
            return text

        # Adiciona um novo slide com título e conteúdo
        def add_slide(title, content):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout de Título e Conteúdo
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]

            set_font(title_shape, bold=True, size=30)
            set_font(title_shape).text = title

            set_font(content_shape).text = content

        # Slide para Nome do Povo, País e Continente
        add_slide("Nome do Povo, País e Continente", 
                  f"{self.data['Nome do povo']}\n{self.data['País']}\n{self.data['Continente']}")

        # Slide para Onde Vivem
        add_slide("Onde Vivem", self.data['Onde vivem'])

        # Slide para População
        add_slide("População", self.data['População'])

        # Slide para Idioma e Tradução
        add_slide("Idioma e Tradução", self.data['Idioma e tradução'])

        # Slide para Religião
        add_slide("Religião", self.data['Religião'])

        # Slide para Relação com o Cristianismo
        add_slide("Relação com o Cristianismo", self.data['Relação com o cristianismo'])

        # Slide para Cristãos e Evangélicos no Brasil e no Mundo
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        set_font(title_shape, bold=True, size=30).text = "Cristãos e Evangélicos"
        content_frame = content_shape.text_frame
        content_frame.clear()  # Limpa o conteúdo padrão

        brasil_text = content_frame.add_paragraph()
        brasil_text.text = "Brasil:\nCristãos no Brasil: {}\nEvangélicos no Brasil: {}".format(
            self.data['Cristãos no Brasil'], self.data['Evangélicos no Brasil'])
        set_font(brasil_text, size=25)

        mundo_text = content_frame.add_paragraph()
        mundo_text.text = "No Mundo:\nCristãos pelo mundo: {}\nEvangélicos pelo mundo: {}".format(
            self.data['Cristãos pelo mundo'], self.data['Evangélicos pelo mundo'])
        set_font(mundo_text, size=25)

        # Slides para Introdução, Como Vivem, Em que Acreditam e Intercessão
        for field_name in ['Introdução', 'Como vivem', 'Em que acreditam', 'Intercessão']:
            add_slide(field_name, self.data[field_name])

        # Adicionar imagem de fundo aos slides
        img_path = self.background_image_path
        if img_path:
            for slide in prs.slides:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = (255, 255, 255)  # Cor de fundo branca
                left = top = 0
                slide.background.fill.user_picture(img_path, left, top, width=None, height=None)

        # Salva a apresentação
        prs.save(os.path.join(self.dest_path, f"{self.file_name}_{date.today()}.pptx"))
